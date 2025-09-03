"""
Multi-format Document Service
Supports PDF, DOCX, PPTX, TXT, RTF, Excel files and more
"""

import pdfplumber # type: ignore
import pytesseract  # type: ignore
from pdf2image import convert_from_bytes # type: ignore
from PIL import Image # type: ignore
from docx import Document as DocxDocument # type: ignore
from pptx import Presentation # type: ignore
from openpyxl import load_workbook # type: ignore
import mammoth # type: ignore
from striprtf.striprtf import rtf_to_text # type: ignore
import chardet # type: ignore
from typing import List, Dict, Any, Optional
import uuid
import io
import os
import mimetypes

class DocumentProcessor:
    """Universal document processor for multiple file formats"""
    
    def __init__(self):
        """Initialize document processor with supported formats"""
        self.supported_formats = {
            '.pdf': self._extract_from_pdf,
            '.docx': self._extract_from_docx,
            '.doc': self._extract_from_doc,
            '.pptx': self._extract_from_pptx,
            '.ppt': self._extract_from_ppt,
            '.xlsx': self._extract_from_excel,
            '.xls': self._extract_from_excel,
            '.txt': self._extract_from_text,
            '.rtf': self._extract_from_rtf,
            '.md': self._extract_from_markdown,
            '.csv': self._extract_from_csv,
        }
    
    def get_supported_formats(self) -> List[str]:
        """Return list of supported file extensions"""
        return list(self.supported_formats.keys())
    
    def is_supported(self, filename: str) -> bool:
        """Check if file format is supported"""
        ext = os.path.splitext(filename.lower())[1]
        return ext in self.supported_formats
    
    def extract_text(self, file, filename: str) -> Dict[str, Any]:
        """Extract text from any supported document format"""
        ext = os.path.splitext(filename.lower())[1]
        
        if not self.is_supported(filename):
            raise ValueError(f"Unsupported file format: {ext}")
        
        try:
            # Reset file pointer
            file.seek(0)
            
            # Get the appropriate extraction function
            extract_func = self.supported_formats[ext]
            
            # Extract text and metadata
            result = extract_func(file, filename)
            
            return {
                'text': result.get('text', ''),
                'metadata': result.get('metadata', {}),
                'format': ext,
                'success': True,
                'error': None
            }
            
        except Exception as e:
            return {
                'text': '',
                'metadata': {},
                'format': ext,
                'success': False,
                'error': str(e)
            }
    
    def _extract_from_pdf(self, file, filename: str) -> Dict[str, Any]:
        """Extract text from PDF with OCR fallback"""
        text = ""
        metadata = {'pages': 0, 'extraction_method': 'text'}
        
        try:
            # First, try to extract text directly
            with pdfplumber.open(file) as pdf:
                metadata['pages'] = len(pdf.pages)
                
                for page_num, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text()
                    if page_text:
                        text += f"[Page {page_num}]\n{page_text}\n\n"
        except Exception as e:
            print(f"Error with pdfplumber extraction: {e}")
        
        # If no text was extracted or very little text, use OCR
        if len(text.strip()) < 100:
            print("Text extraction yielded little content. Attempting OCR...")
            try:
                file.seek(0)
                file_bytes = file.read()
                
                # Convert PDF pages to images
                images = convert_from_bytes(file_bytes, dpi=300, fmt='jpeg')
                metadata['extraction_method'] = 'ocr'
                
                ocr_text = ""
                for i, image in enumerate(images, 1):
                    print(f"Processing page {i} with OCR...")
                    page_text = pytesseract.image_to_string(image, lang='eng')
                    if page_text.strip():
                        ocr_text += f"[Page {i}]\n{page_text.strip()}\n\n"
                
                if ocr_text.strip():
                    text = ocr_text
                    
            except Exception as e:
                print(f"Error with OCR extraction: {e}")
        
        return {'text': text.strip(), 'metadata': metadata}
    
    def _extract_from_docx(self, file, filename: str) -> Dict[str, Any]:
        """Extract text from DOCX files"""
        try:
            # Reset file pointer
            file.seek(0)
            
            # Check if file is large enough to be a valid DOCX
            file_content = file.read()
            if len(file_content) < 100:
                raise Exception("File too small to be a valid DOCX document")
            
            # Create a BytesIO object from the content for python-docx compatibility
            import io
            file_bytes = io.BytesIO(file_content)
            
            doc = DocxDocument(file_bytes)
            
            text_parts = []
            metadata = {
                'paragraphs': len(doc.paragraphs),
                'tables': len(doc.tables),
                'sections': len(doc.sections)
            }
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text.strip())
            
            # Extract text from tables
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_text.append(" | ".join(row_text))
                
                if table_text:
                    text_parts.append("\\n".join(table_text))
            
            # Try to get document properties
            try:
                core_props = doc.core_properties
                if core_props.title:
                    metadata['title'] = core_props.title
                if core_props.author:
                    metadata['author'] = core_props.author
                if core_props.subject:
                    metadata['subject'] = core_props.subject
            except:
                pass
                
            return {'text': '\\n\\n'.join(text_parts), 'metadata': metadata}
            
        except Exception as e:
            raise Exception(f"Error extracting DOCX: {e}")
    
    def _extract_from_doc(self, file, filename: str) -> Dict[str, Any]:
        """Extract text from legacy DOC files using mammoth"""
        try:
            file.seek(0)
            file_content = file.read()
            
            # Create BytesIO object for mammoth compatibility
            import io
            file_bytes = io.BytesIO(file_content)
            
            result = mammoth.extract_raw_text(file_bytes)
            
            metadata = {
                'extraction_method': 'mammoth',
                'warnings': result.messages if hasattr(result, 'messages') else []
            }
            
            return {'text': result.value, 'metadata': metadata}
            
        except Exception as e:
            raise Exception(f"Error extracting DOC: {e}")
    
    def _extract_from_pptx(self, file, filename: str) -> Dict[str, Any]:
        """Extract text from PowerPoint PPTX files"""
        try:
            file.seek(0)
            file_content = file.read()
            
            # Create BytesIO object for python-pptx compatibility
            import io
            file_bytes = io.BytesIO(file_content)
            
            prs = Presentation(file_bytes)
            
            text_parts = []
            metadata = {
                'slides': len(prs.slides),
                'slide_layouts': len(prs.slide_layouts)
            }
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"[Slide {slide_num}]\\n"
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += f"{shape.text.strip()}\\n"
                    
                    # Handle tables in slides
                    if shape.shape_type == 19:  # Table
                        try:
                            table = shape.table
                            for row in table.rows:
                                row_text = []
                                for cell in row.cells:
                                    if cell.text.strip():
                                        row_text.append(cell.text.strip())
                                if row_text:
                                    slide_text += " | ".join(row_text) + "\\n"
                        except:
                            pass
                
                if slide_text.strip() != f"[Slide {slide_num}]":
                    text_parts.append(slide_text)
            
            return {'text': '\\n\\n'.join(text_parts), 'metadata': metadata}
            
        except Exception as e:
            raise Exception(f"Error extracting PPTX: {e}")
    
    def _extract_from_ppt(self, file, filename: str) -> Dict[str, Any]:
        """Extract text from legacy PPT files (limited support)"""
        return {
            'text': f"Legacy PPT format detected for {filename}. Please convert to PPTX for better text extraction.",
            'metadata': {'note': 'Legacy PPT format - limited support'}
        }
    
    def _extract_from_excel(self, file, filename: str) -> Dict[str, Any]:
        """Extract text from Excel files (XLSX/XLS)"""
        try:
            file.seek(0)
            file_content = file.read()
            
            # Create BytesIO object for openpyxl compatibility
            import io
            file_bytes = io.BytesIO(file_content)
            
            workbook = load_workbook(file_bytes, read_only=True)
            
            text_parts = []
            metadata = {
                'worksheets': len(workbook.worksheets),
                'sheet_names': workbook.sheetnames
            }
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = f"[Sheet: {sheet_name}]\\n"
                
                # Find the data range
                max_row = sheet.max_row
                max_col = sheet.max_column
                
                if max_row > 0 and max_col > 0:
                    for row in sheet.iter_rows(min_row=1, max_row=min(max_row, 100), values_only=True):
                        row_text = []
                        for cell_value in row:
                            if cell_value is not None and str(cell_value).strip():
                                row_text.append(str(cell_value).strip())
                        
                        if row_text:
                            sheet_text += " | ".join(row_text) + "\\n"
                
                if sheet_text.strip() != f"[Sheet: {sheet_name}]":
                    text_parts.append(sheet_text)
            
            return {'text': '\\n\\n'.join(text_parts), 'metadata': metadata}
            
        except Exception as e:
            raise Exception(f"Error extracting Excel: {e}")
    
    def _extract_from_text(self, file, filename: str) -> Dict[str, Any]:
        """Extract text from plain text files with encoding detection"""
        try:
            file.seek(0)
            raw_data = file.read()
            
            # Detect encoding
            encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'
            
            # Decode text
            if isinstance(raw_data, bytes):
                text = raw_data.decode(encoding, errors='replace')
            else:
                text = str(raw_data)
            
            metadata = {
                'encoding': encoding,
                'lines': len(text.split('\\n')),
                'characters': len(text)
            }
            
            return {'text': text, 'metadata': metadata}
            
        except Exception as e:
            raise Exception(f"Error extracting text: {e}")
    
    def _extract_from_rtf(self, file, filename: str) -> Dict[str, Any]:
        """Extract text from RTF files"""
        try:
            file.seek(0)
            rtf_content = file.read()
            
            if isinstance(rtf_content, bytes):
                rtf_content = rtf_content.decode('utf-8', errors='replace')
            
            text = rtf_to_text(rtf_content)
            
            metadata = {
                'original_format': 'rtf',
                'characters': len(text)
            }
            
            return {'text': text, 'metadata': metadata}
            
        except Exception as e:
            raise Exception(f"Error extracting RTF: {e}")
    
    def _extract_from_markdown(self, file, filename: str) -> Dict[str, Any]:
        """Extract text from Markdown files"""
        try:
            file.seek(0)
            raw_data = file.read()
            
            # Detect encoding
            encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'
            
            if isinstance(raw_data, bytes):
                text = raw_data.decode(encoding, errors='replace')
            else:
                text = str(raw_data)
            
            metadata = {
                'format': 'markdown',
                'encoding': encoding,
                'lines': len(text.split('\\n'))
            }
            
            return {'text': text, 'metadata': metadata}
            
        except Exception as e:
            raise Exception(f"Error extracting Markdown: {e}")
    
    def _extract_from_csv(self, file, filename: str) -> Dict[str, Any]:
        """Extract text from CSV files"""
        try:
            import csv
            
            file.seek(0)
            raw_data = file.read()
            
            # Detect encoding
            encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'
            
            if isinstance(raw_data, bytes):
                text_data = raw_data.decode(encoding, errors='replace')
            else:
                text_data = str(raw_data)
            
            # Parse CSV
            csv_reader = csv.reader(io.StringIO(text_data))
            
            rows = []
            for row_num, row in enumerate(csv_reader, 1):
                if row:  # Skip empty rows
                    rows.append(" | ".join([cell.strip() for cell in row if cell.strip()]))
                    if row_num > 500:  # Limit to first 500 rows
                        break
            
            text = '\\n'.join(rows)
            
            metadata = {
                'format': 'csv',
                'encoding': encoding,
                'rows_processed': len(rows)
            }
            
            return {'text': text, 'metadata': metadata}
            
        except Exception as e:
            raise Exception(f"Error extracting CSV: {e}")


# Legacy functions for backward compatibility
def extract_text_from_pdf(file) -> str:
    """Legacy function - use DocumentProcessor instead"""
    processor = DocumentProcessor()
    result = processor._extract_from_pdf(file, "document.pdf")
    return result['text']

def chunk_text(text: str, chunk_size: int = 800, overlap: int = 100) -> List[str]:
    """Split text into overlapping chunks with better context preservation"""
    
    # First, try to split by paragraphs
    paragraphs = text.split('\\n\\n')
    
    chunks = []
    current_chunk = ""
    
    for paragraph in paragraphs:
        paragraph = paragraph.strip()
        if not paragraph:
            continue
            
        # If adding this paragraph would exceed chunk size, save current chunk and start new one
        if len(current_chunk) + len(paragraph) + 2 > chunk_size and current_chunk:
            if current_chunk.strip():
                chunks.append(current_chunk.strip())
            
            # Start new chunk with overlap from previous chunk
            if overlap > 0 and current_chunk:
                words = current_chunk.split()
                overlap_words = words[-overlap//4:] if len(words) > overlap//4 else words[-10:]
                current_chunk = " ".join(overlap_words) + " " + paragraph
            else:
                current_chunk = paragraph
        else:
            # Add paragraph to current chunk
            if current_chunk:
                current_chunk += "\\n\\n" + paragraph
            else:
                current_chunk = paragraph
    
    # Add the last chunk
    if current_chunk.strip():
        chunks.append(current_chunk.strip())
    
    # If we ended up with very few chunks, fall back to word-based chunking
    if len(chunks) < 3:
        words = text.split()
        chunks = []
        
        for i in range(0, len(words), chunk_size - overlap):
            chunk = " ".join(words[i:i + chunk_size])
            if chunk.strip():
                chunks.append(chunk.strip())
    
    return chunks

def generate_doc_id() -> str:
    """Generate unique document ID"""
    return str(uuid.uuid4())

# Global document processor instance
document_processor = DocumentProcessor()
